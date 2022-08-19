import pydash
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import ImageColor
from tpip_pptx.constants import CommandRegexSub,CommandRegex
from tpip_pptx.tag import Tag

class Text(Tag):
    def __init__(self):
        pass

    def text_replace(self,commands_dic,presentation,slide,shape,slides,dataObj):
        pattern = CommandRegex.TEXT.value
        matching_val = super().get_object_values(pattern, shape, dataObj)
       
        for val in matching_val:
            if type(matching_val[val]) == dict:
                super().replace_tags(str(f"{CommandRegexSub.INS.value} {val} +++"), str(matching_val[val].get('text')), shape)
            else:
                super().replace_tags(str(f"{CommandRegexSub.INS.value} {val} +++"), str(matching_val[val]), shape)

            if type(matching_val[val]) == dict:
                text = matching_val[val].get('text')
                styles = matching_val[val].get('styles')

                font_color = pydash.get(styles, "font_color")
                styles['font_color'] = ImageColor.getcolor(font_color, "RGB")                       

                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:             
                            run.text = text
                            if 'name' in styles:
                                run.font.name = styles['name'] 
                            if 'size' in styles:                       
                                run.font.size = Pt(int(styles['size'] ))
                            if 'bold' in styles:                       
                                run.font.bold = bool(styles['bold'] )
                            if 'italic' in styles:
                                run.font.italic = bool(styles['italic'])
                            if 'underline' in styles:
                                run.font.underline = bool(styles['underline'])
                            if 'font_color' in styles:
                                run.font.color.rgb = RGBColor(styles["font_color"][0], styles["font_color"][1], styles["font_color"][2])
                            

    