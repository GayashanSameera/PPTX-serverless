import pydash

from tpip_pptx.tag import Tag
from copy import deepcopy
from pptx.table import _Cell
from tpip_pptx.constants import CommandRegex, CommandRegexSub

class Toc(Tag):
    def __init__(self):
        pass 

    def drow_toc(self,presentation,data):
        ids = self.identify_tags_with_page_number(presentation)
        slides = [slide for slide in presentation.slides]
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if(str(f"{CommandRegexSub.TOC.value}") in shape.text):
                        pattern = CommandRegex.TOC.value
                        matches = super().get_tag_content(pattern, shape)
                        data_path = matches[0]
                        datam = pydash.get(data, data_path)
                        self.update_toc_table(slide,datam,ids)
                        super().replace_tags(str(f"{CommandRegexSub.TOC.value} {data_path} +++"), "", shape)
                        return

    def update_toc_table(self,slide,datam,ids):
        for shape in slide.shapes:
            if shape.has_table:
                self.execute_table_drower(shape.table, datam, ids)

    def execute_table_drower(self,table, data,ids):
        row_index = 0
        for row in data:
            if row_index > 0:
                self.add_new_row_to_existing_table(table)
            
            cell_1 = table.cell(row_index, 0)
            cell_1.text = row["text"]
            cell_2 = table.cell(row_index, 2)
            row_id = row["id"]
            cell_2.text = str(ids[row_id])

            row_index += 1

    def add_new_row_to_existing_table(self,table):
        new_row = deepcopy(table._tbl.tr_lst[0])
        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''

            table._tbl.append(new_row) 
            return table.rows[0]
    

    def identify_tags_with_page_number(self,presentation):
        res = {}
        slides = [slide for slide in presentation.slides]
        pattern = CommandRegex.TOC_IDS.value
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    matches = super().get_tag_content(pattern, shape)
                    if(matches and len(matches) > 0):
                        match_string = matches[0]
                        page_number = slides.index(slide) + 1
                        if "," in match_string:
                            id_array = match_string.split(",")
                            for id in id_array:
                                res[id] = page_number
                        else:
                            res[match_string] = page_number

                super().replace_tags(str(f"{CommandRegexSub.TOC_IDS.value} {matches} +++"), "", shape)
        return res