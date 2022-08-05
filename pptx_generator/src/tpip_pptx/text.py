from constants import CommandRegexSub
from tag import Tag
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from PIL import ImageColor
import pydash
from pptx.util import Inches
import enum
import re

class Text(Tag):
    def __init__(self):
        pass

    def text_replace(self, slide, pattern, shape):
        match, object_value = super().get_object_values(pattern, shape)
        super().replace_tags(str(f"{CommandRegexSub.INS.value} {match} +++"), str(object_value.get('text')), shape)

        if type(object_value) == dict:
            text = object_value.get('text')
            styles = object_value.get('styles')

            font_color = pydash.get(styles, "font_color")
            styles['font_color'] = ImageColor.getcolor(font_color, "RGB")                       

            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:             
                        run.text = text
                        if 'name' in styles:
                            run.font.name = styles['name'] 
                        if 'size' in styles:                       
                            run.font.size = Pt(int(styles['size'] ))
                        if 'bold' in styles:                       
                            run.font.bold = bool(styles['bold'] )
                        if 'italic' in styles:
                            run.font.italic = bool(styles['italic'])
                        if 'underline' in styles:
                            run.font.underline = bool(styles['underline'])
                        if 'font_color' in styles:
                            run.font.color.rgb = RGBColor(styles["font_color"][0], styles["font_color"][1], styles["font_color"][2])
                        

    def text_tag_update(pattern, text):
        match, object_value = super().get_object_values_string(pattern, text)
        if (object_value != False):
            current_text = current_text.replace(str(f"{CommandRegexSub.INS.value} {match} +++"), str(object_value))