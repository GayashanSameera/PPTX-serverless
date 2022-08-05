import pydash
import math
from pptx.oxml.xmlchemy import OxmlElement
from copy import deepcopy
from pptx.table import Table, _Row, _Column, _Cell
from pptx.dml.color import RGBColor
from pptx.util import Inches

from tag import Tag
from constants import CommandRegexSub,CommandRegex


class Table(Tag):
        def __init__(self):
            pass 

        def update_table_text(self,commands_dic,presentation, slide, shape, slide_index, dataObj):
            pattern = CommandRegex.UPDATE_TABLE_TEXT.value[0]
            matches = super().get_tag_content(pattern, shape)

            if( not matches or len(matches) < 1):
                return

            for match in matches:
                styles = False
                data = False
                data_path = False
                table_id = match

                if "DATA" in match:
                    data_path = match.split(" DATA ")[1]
                    table_id = match.split(" DATA ")[0]

                if data_path and data_path in dataObj:
                    data = dataObj[data_path]

                if data and "styles" in data:
                    styles = data["styles"]

                table_id_tag = str(f"{CommandRegexSub.TB_ID.value} {table_id} +++")
                for _shape in slide.shapes:
                    if _shape.has_table: 
                        for row in _shape.table.rows:
                            for cell in row.cells:
                                if table_id_tag in cell.text:
                                    self.execute_table_tags(_shape, _shape.table, data, styles)
                                    new_text = cell.text.replace(str(f"{CommandRegexSub.TB_ID.value} {table_id} +++"), "")
                                    cell.text = new_text
                                    break

                super().replace_tags(str(f"{CommandRegexSub.TB_TX_UP.value} {match} +++"), "", shape)

        def execute_table_tags(self,shape , table, dataObj, styles):
            row_index = 0
            for row in table.rows:
                col_index = 0
                for cell in row.cells:
                    pattern_for = CommandRegex.PATTERN_FOR.value[0]
                    matches_for = super().get_tag_from_string(pattern_for, cell.text)
                    if( matches_for and len(matches_for) > 0):
                        for match in matches_for:
                            pattern_condition = CommandRegex.PATTERN_CONDITION.value[0]
                            matched_condition = super().get_tag_from_string(pattern_condition,match)

                            pattern_content = CommandRegex.PATTERN_CONTENT.value[0]
                            matched_content = super().get_tag_from_string(pattern_content,match)
                            for contidion in matched_condition:
                                object_value = pydash.get(dataObj, contidion)
                                text_result = ""
                                if(object_value):
                                    data_count = 1
                                    for _data in object_value:
                                        updated_data = super().text_tag_update(matched_content[0],_data)
                                        if(updated_data and updated_data["text"] and len(object_value) > data_count):
                                            text_result += updated_data["text"] + "\n"
                                        elif(updated_data and updated_data["text"]):
                                            text_result += updated_data["text"]
                                        data_count += 1
                                new_text = cell.text.replace(str(f"{CommandRegexSub.FOR.value} {match} {CommandRegexSub.FOR_END.value}"), text_result)
                                cell.text = new_text
                                try:
                                    self.table_styles(cell,row_index,col_index,styles)
                                except ValueError:
                                    print("error")
                        

                    pattern_text = CommandRegex.TEXT.value[0]
                    matches_text_update = super().get_tag_from_string(pattern_text, cell.text)
                    if( matches_text_update and len(matches_text_update) > 0):
                        for match in matches_text_update:
                            new_text = cell.text.replace(str(f"{CommandRegexSub.INS.value} {match} +++"), pydash.get(dataObj, match))
                            cell.text = new_text
                            try:
                                self.table_styles(cell,row_index,col_index,styles)
                            except ValueError:
                                print("error")
                    col_index += 1 
                row_index +=1     
            
        def table_styles(self,cell,row_index,col_index,styles):
            try:
                row_st_index = str(f'rw_{row_index}')
                col_st_index = str(f'cl_{col_index}')
                para_index = 0
                for paragraph in cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[para_index]
                    if(styles and 'all' in styles):
                        common_styles = styles['all']

                        if('font_size' in common_styles):
                            para.font.size = Pt(common_styles['font_size'])
                        if('font_name' in common_styles):
                            para.font.name = common_styles['font_name']
                        if('bold' in common_styles):
                            para.font.bold = common_styles['bold']
                        if('italic' in common_styles):
                            para.font.italic = common_styles['italic']
                        if("font_color" in common_styles):
                            para.font.color.rgb = RGBColor(common_styles["font_color"][0], common_styles["font_color"][1],common_styles["font_color"][2])
                        if("alignment" in common_styles):
                            if common_styles["alignment"] == "center":
                                para.alignment = PP_ALIGN.CENTER
                        if("background_color" in common_styles):
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(common_styles["background_color"][0], common_styles["background_color"][1],common_styles["background_color"][2])

                    if(styles and row_st_index in styles):
                        _styles = styles[row_st_index]
                    if( "column_indexes" in _styles):
                        if(col_index in _styles["column_indexes"]):
                            if('font_size' in _styles):
                                para.font.size = Pt(_styles['font_size'])
                            if('font_name' in _styles):
                                para.font.name = _styles['font_name']
                            if('bold' in _styles):
                                para.font.bold = _styles['bold']
                            if('italic' in _styles):
                                para.font.italic = _styles['italic']
                            if("font_color" in _styles):
                                para.font.color.rgb = RGBColor(_styles["font_color"][0], _styles["font_color"][1],_styles["font_color"][2])
                            if("background_color" in _styles):
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(_styles["background_color"][0], _styles["background_color"][1],_styles["background_color"][2])
                    else:
                        if('font_size' in _styles):
                            para.font.size = Pt(_styles['font_size'])
                        if('font_name' in _styles):
                            para.font.name = _styles['font_name']
                        if('bold' in _styles):
                            para.font.bold = _styles['bold']
                        if('italic' in _styles):
                            para.font.italic = _styles['italic']
                        if("font_color" in _styles):
                            para.font.color.rgb = RGBColor(_styles["font_color"][0], _styles["font_color"][1],_styles["font_color"][2])
                        if("background_color" in _styles):
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(_styles["background_color"][0], _styles["background_color"][1],_styles["background_color"][2])

                if(styles and col_st_index in styles):
                    col_styles = styles[col_st_index]

                    if('font_size' in col_styles):
                        para.font.size = Pt(col_styles['font_size'])
                    if('font_name' in col_styles):
                        para.font.name = col_styles['font_name']
                    if('bold' in col_styles):
                        para.font.bold = col_styles['bold']
                    if('italic' in col_styles):
                        para.font.italic = col_styles['italic']
                    if("font_color" in col_styles):
                        para.font.color.rgb = RGBColor(col_styles["font_color"][0], col_styles["font_color"][1],col_styles["font_color"][2])
                    if("background_color" in col_styles):
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(col_styles["background_color"][0], col_styles["background_color"][1],col_styles["background_color"][2])
        
                para_index += 1
            except ValueError:
                print("error")
                
        
        def replace_tables(self,pattern,presentation,slide,shape,slide_index,dataObj):
            match , object_value = super().get_object_values(pattern, shape)
            if(object_value):
                super().replace_tags(str(f"{CommandRegexSub.TB_ADD.value} {match} +++"), "", shape)
                self.create_table(presentation, slide, shape, slide_index, object_value)
                
                
        def create_table(presentation, slide, shape, slide_index, dataObj):
            row_count = pydash.get(dataObj,dataObj.cashFlows.row_count,default=5)
            cols = pydash.get(dataObj,dataObj.cashFlows.colum_count,default=3)
            headers = pydash.get(dataObj,dataObj.cashFlows.headers)
            row_data = pydash.get(dataObj,dataObj.cashFlows.rows)
            styles = pydash.get(dataObj,dataObj.cashFlows.styles)
            table_count_per_slide = pydash.get(dataObj,dataObj.cashFlows.table_count_per_slide,default=2)
            
            
            total_rows = len(row_data)
            total_table_count = math.ceil(total_rows / row_count)

            slide_count = math.ceil(total_table_count / table_count_per_slide) 
            extra_slide_count = slide_count - 1

            s = 0
            end = 0
            current_slide = 0
            while s < slide_count:
                j = 0
                left = 1

                if total_table_count > table_count_per_slide :
                    table_count = table_count_per_slide
                else:
                    table_count = total_table_count

                slide_row_start = 0
                slide_row_end = 0
                extra_slide_exists = False

                if(s > 0 and (not presentation.slides[slide_index + s])):
                    break

                if(s > 0 and current_slide != s):
                    for new_slide_shape in presentation.slides[slide_index + s].shapes:
                        extra_slide_exists = is_extra_slide(presentation, slide_index + s, True)
                        if(extra_slide_exists):
                            break
                    current_slide = s

                if(s > 0 and (not extra_slide_exists)):
                    return 
                

                while j < table_count:
                    shape = presentation.slides[slide_index + s].shapes.add_table(row_count + 1, cols, Inches(left) , Inches(styles["top"]), Inches(styles["width"]), Inches(styles["row_height"]))
                    table = shape.table
                    tables_headers(table, headers)

                    slide_row_start = (row_count * j ) + 1 + end
                    slide_row_end = slide_row_start + (row_count - 1)

                    tables_rows(table,row_data, slide_row_start, slide_row_end ,total_rows )
                    left += 2
                    j += 1

            end = slide_row_end
            total_table_count -= table_count_per_slide
            s += 1
            
            def is_extra_slide(presentation, slide_index, remove_tag):
                extra_slide_exists = False
                extra = 'EXTRA_SLIDE'
                for new_slide_shape in presentation.slides[slide_index].shapes:
                    if new_slide_shape.has_text_frame:
                        matches = super().check_tag_exist(extra, new_slide_shape)
                        if(matches):
                            extra_slide_exists = True
                            if(remove_tag):
                                super().replace_tags(str(f"+++ {extra} +++"), "", new_slide_shape)
                            break
                return extra_slide_exists
            
            
            def tables_headers(table, headers):
                i = 0
                for header in headers:
                    table.cell(0, i).text = header
                    cell = table.cell(0, i)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(200, 253, 251)
                    _set_cell_border(cell,"949595", '12000')
                    para = cell.text_frame.paragraphs[0]
                    para.font.bold = True
                    para.font.size = Pt(5)
                    para.font.name = 'Comic Sans MS'
                    para.font.color.rgb = RGBColor(19, 170, 246)
                    table.columns[i].width = Inches(0.6)
                    i += 1
                    
            def tables_rows(table, rowData, start, end,totalRows):
                j = start
                cell_start = 1

                entCount = end
                if end > totalRows:
                    entCount = totalRows

                while j < entCount + 1:
                    element = rowData[j - 1]
                    k = 0
                    if element:
                        for key, value in element.items():
                            table.cell(cell_start, k).text = value
                            cell = table.cell(cell_start, k)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(228, 228, 228)
                            _set_cell_border(cell,"949595", '12000')
                            para = cell.text_frame.paragraphs[0]
                            para.font.size = Pt(6)
                            para.font.name = 'Comic Sans MS'
                            k += 1
                    j += 1
                    cell_start += 1
                    
            def SubElement(parent, tagname, **kwargs):
                element = OxmlElement(tagname)
                element.attrib.update(kwargs)
                parent.append(element)
                return element
            
            def _set_cell_border(cell, border_color="000000", border_width='12700'):
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
                    ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
                    solidFill = SubElement(ln, 'a:solidFill')
                    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
                    prstDash = SubElement(ln, 'a:prstDash', val='solid')
                    round_ = SubElement(ln, 'a:round')
                    headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
                    tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
                    
                    
        def drow_tables(self,presentation, slide, shape, slide_index, dataObj):
            pattern = CommandRegex.TABLE_DRAW.value[0]
            matches = super().get_tag_content(pattern, shape)
            if( not matches or len(matches) < 1):
                return

            for match in matches:
                styles = False
                data = False
                data_path = False
                table_id = match

            if "DATA" in match:
                data_path = match.split(" DATA ")[1]
                table_id = match.split(" DATA ")[0]

            if data_path and data_path in dataObj:
                data = dataObj[data_path]

            if data and "styles" in data:
                styles = data["styles"]

            table_id_tag = str(f"{CommandRegexSub.TB_ID.value} {table_id} +++")
            for _shape in slide.shapes:
                if _shape.has_table: 
                    for row in _shape.table.rows:
                        for cell in row.cells:
                            if table_id_tag in cell.text:
                                self.execute_table_drower(_shape.table, data, styles)
                                new_text = cell.text.replace(str(f"{CommandRegexSub.TB_ID.value} {table_id} +++"), "")
                                cell.text = new_text
                                break

            super().replace_tags(str(f"{CommandRegexSub.TB_DRW.value} {match} +++"), "", shape)
        
        def execute_table_drower(self,table, dataObj,styles):
            row_data = pydash.get(dataObj,dataObj.cashFlows.rows,default=[])
            row_index = 1
            for row in row_data:
                colum_index = 0
                if row_index > 1:
                    self.add_new_row_to_existing_table(table)
                for column in row:
                    cell = table.cell(row_index, colum_index)
                    cell.text = column
                    
                    try:
                        self.table_styles(cell,row_index,colum_index,styles)
                    except ValueError:
                        print("error")

                    colum_index += 1
                row_index += 1
                
        def add_new_row_to_existing_table(table):
            new_row = deepcopy(table._tbl.tr_lst[1])
            for tc in new_row.tc_lst:
                cell = _Cell(tc, new_row.tc_lst)
                cell.text = '' # defaulting cell contents to empty text

                table._tbl.append(new_row) 
                return table.rows[1]