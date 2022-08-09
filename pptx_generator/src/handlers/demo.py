from ctypes import alignment
import pydash
from pptx.util import Inches
from pptx import Presentation
import enum
import re
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from PIL import ImageColor

# from pptx_generator.tpip_pptx.constants import CommandRegex, CommandRegexSub, Command
from lxml import etree



dataObj = {
       
        "schemeName": {
            "text": "XYZ Pension Scheme",
            "styles": {
                    "name": "Comic Sans MS",
                    "size": "40",
                    "italic": "True",
                    "bold": "True",
                    "alignment": "center",
                    "underline": "True",
                    "font_color": "#FFFF00"
                
            }
        },
        "title": {
            "text": "Q2 2021 Summary Reporttt",
            "styles": {
                    "name": "Comic Sans MS",
                    "size": "50",
                    "italic": "True",
                    "bold": "True",
                    "alignment": "center",
                    "underline": "True",
                    "font_color": "#FF5733"
                }
            },
        "heading": {
            "text": "Investment performance to 30 June 2021",
            },
        'table_name': "Sample table to delete",
        "remove_table_1": True,
        'table_name_row': "Sample table to delete row",
        "table_1_row_3_present": False,
        'table_name_column': "Sample table to delete column",
        "table_1_col_4_present": False,
	
    }

prs = Presentation('demo.pptx')
print(prs)


    
class CommandRegex(enum.Enum):
    IMAGE = r'\+\+\+IM (.*?) \+\+\+',
    TEXT = r'\+\+\+INS (.*?) \+\+\+'
    UPDATE_TABLE_TEXT = r'\+\+\+TB_TX_UP (.*?) \+\+\+',
    CREATE_TABLE = r'\+\+\+TB_ADD (.*?) \+\+\+',
    PATTERN_IF = r'\+\+\+IF (.*?)IF-END\+\+\+'
    PATTERN_FOR = r'\+\+\+FOR (.*?) FOR-END\+\+\+',
    PATTERN_CONTENT = r'\<\<(.*?)\>\>',
    PATTERN_CONDITION = r'\(\((.*?)\)\)',
    TABLE_DRAW = r'\+\+\+TB_DRW (.*?) \+\+\+'
    TABLE_REMOVE = r'\+\+\+TABLE_REMOVE (.*?) \+\+\+'
    TABLE_ROW_REMOVE = r'\+\+\+TABLE_ROW_REMOVE (.*?) \+\+\+'
    TABLE_COLUMN_REMOVE = r'\+\+\+TABLE_COLUMN_REMOVE (.*?) \+\+\+'
    
class Command(enum.Enum):
        IF_CONDITION = "if_condition"
        FOR_LOOP = "for_loop"
        REPLACE_IMAGE = "replace_images"
        REPLACE_TABLE = "replace_table"
        UPDATE_TABLE_TEXT = "update_table_text"
        DRAW_TABLE = "draw_tables"
        TEXT_REPLACE = "text_replace"
        REMOVE_TABLE = "remove_tables"

class CommandRegexSub(enum.Enum):
    IMG = '+++IM'
    INS = '+++INS'
    TB_ADD = '+++TB_ADD'
    TB_TX_UP = '+++TB_TX_UP'
    TB_ID ='+++TB_ID'
    TB_DRW = '+++TB_DRW'
    FOR = '+++FOR'
    FOR_END = 'FOR-END+++'
    IF = '+++IF'
    IF_END = 'IF-END+++'
    RW_ID ='+++RW_ID'  

class Tag:
        def __init__(self,pattern):
                self.pattern = pattern
        
        def get_tag_content(self,pattern, shape):
             print('inside matches', pattern)   
           
             print('shape.text',shape.text)
             matches = re.findall(pattern,shape.text)
             print('matches', matches)
             return matches
        
        def replace_tags(self,replaced_for,replaced_text,shape):
            print("vvv", shape)
            # print('replace_for',replaced_for)
            # print('replace_text',replaced_text)
            if shape.has_text_frame:
                text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:             
                    cur_text = run.text
                    print('cur_text',cur_text)
                    new_text = cur_text.replace(replaced_for, replaced_text)
                    print('new_text',new_text)
                    run.text = new_text
                    
        def get_object_values(self,pattern,shape):
            # print('seee', shape)
            matches = self.get_tag_content(pattern,shape)
            # print('Matches', matches)a
            if( not matches or len(matches) < 1):
                return
            for match in matches:
                object_value = pydash.get(dataObj, match)
                return match , object_value
              
        def get_object_values_string(self,pattern,shape):
             if( not matches or len(matches) < 1):
               return { "text": text }

             for match in matches:
               object_value = pydash.get(dataObj, match, False)
               if(object_value != False):
                    current_text = current_text.replace(str(f"+++INS {match} +++"), str(object_value))

    
             return { "text": current_text }
                
            
        def get_tag_from_string(self,pattern,string):
            print(" pattern string", pattern)
            matches = re.findall(pattern, string)
            return matches
        
        def eval_executor(self,logic,dataObj):
            return eval(logic,dataObj)
        
        def text_tag_update(self,pattern, text,dataObj):
            match, object_value = super().get_object_values_string(pattern, text,dataObj)
            if (object_value != False):
                current_text = current_text.replace(str(f"{CommandRegexSub.INS.value} {match} +++"), str(object_value))

class Image(Tag):
        def __init__(self):
            pass
            
        def replace_images(self,slide,pattern,shape): 
          
          match , object_value = super().get_object_values(pattern,shape)
          
          url = pydash.get(object_value, "url")
          left = pydash.get(object_value, "size.left")
          height = pydash.get(object_value, "size.height")
          top = pydash.get(object_value, "size.top")
          width = pydash.get(object_value, "size.width")
          
          
          slide.shapes.add_picture(url, Inches(left), Inches(top), Inches(width) ,Inches(height) )
          super().replace_tags(str(f"+++IM {match} +++"),"", shape)
          
        def printHello(self,name):
           print(name)
        
 
class Text(Tag):
  
    def __init__(self):
        pass

    def add_styles(self, shape):
        print("shape", shape)
          
    def text_replace(self, slide, pattern,shape):
        match, object_value =super().get_object_values(pattern,shape)
        super().replace_tags(str(f"+++INS {match} +++"), str(object_value.get('text')), shape)

        print("Type ", type(object_value))
        if type(object_value) == dict:
            text = object_value.get('text')
            styles = object_value.get('styles')

            font_color = pydash.get(styles, "font_color")
            styles['font_color'] = ImageColor.getcolor(font_color, "RGB")
            background_color = pydash.get(styles, "background_color")
                       

            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:             
                        run.text = text
                        if 'name' in styles:
                            run.font.name = styles['name'] 
                        if 'size' in styles:                       
                            run.font.size = Pt(int(styles['size'] ))
                        if 'bold' in styles:                       
                            run.font.bold = bool(styles['bold'] )
                        if 'italic' in styles:
                            run.font.italic = bool(styles['italic'])
                        if 'underline' in styles:
                            run.font.underline = bool(styles['underline'])
                        if 'font_color' in styles:
                            run.font.color.rgb = RGBColor(styles["font_color"][0], styles["font_color"][1], styles["font_color"][2])
      

class Table(Tag):
    def __init__(self):
        pass

    def remove_row(self,table, row):
        tbl = table._tbl
        tr = row._tr
        tbl.remove(tr)

    #remove tables
    def remove_tables(self,slide,content):
        table_remove_pattern = CommandRegex.TABLE_REMOVE.value
        table_remove_matches = super().get_tag_from_string(table_remove_pattern, content)
        
        if( table_remove_matches and len(table_remove_matches) > 0):
            table_remove_index_matches = table_remove_matches[0]
            table_id_tag = str(f"{CommandRegexSub.TB_ID.value} {table_remove_index_matches} +++")
            _shap_count = 0
            for _shape in slide.shapes:
                if _shape.has_table: 
                    for row in _shape.table.rows:
                        for cell in row.cells:
                            if table_id_tag in cell.text:
                                old_picture = slide.shapes[_shap_count]
                                old_pic = old_picture._element
                                old_pic.getparent().remove(old_pic)
                                break
                _shap_count += 1

    def remove_table_rows(self,slide,content):
        print("Inside")
        table_row_remove_pattern = CommandRegex.TABLE_ROW_REMOVE.value
        table_row_remove_matches = super().get_tag_from_string(table_row_remove_pattern, content)
        print("table row remove matches", table_row_remove_matches)
        if( table_row_remove_matches and len(table_row_remove_matches) > 0):
            print("IF")
            table_row_remove_index_matches = table_row_remove_matches[0]
            print("table_row_remove_index_matches",table_row_remove_index_matches)
            table_rw_id_tag = str(f"{CommandRegexSub.RW_ID.value} {table_row_remove_index_matches} +++")
            for _shape in slide.shapes:
                if _shape.has_table: 
                    for row_idx, row in enumerate(_shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if table_row_remove_index_matches in cell.text:
                                row_deleted = _shape.table.rows[row_idx]
                                Table.remove_row(_shape.table, row_deleted)
                                break

    def remove_table_column(self,slide,content):
        table_column_remove_pattern = CommandRegex.TABLE_COLUMN_REMOVE.value[0]
        table_column_remove_matches = super().get_tag_from_string(table_column_remove_pattern, content)
        if( table_column_remove_matches and len(table_column_remove_matches) > 0):
            table_column_remove_index_matches = table_column_remove_matches[0]
            for _shape in slide.shapes:
                if _shape.has_table:
                    colum_index = ""
                    for row_idx, row in enumerate(_shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if table_column_remove_index_matches in cell.text:
                                colum_index = col_idx
                                break

                    for row_idx, row in enumerate(_shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if col_idx == colum_index:
                                cell._tc.delete()

                    tree = etree.ElementTree(_shape.table._tbl)
                    for e in tree.iter():
                        if(tree.getpath(e) == tree.getpath(_shape.table.columns[colum_index]._gridCol)):
                            e.getparent().remove(e)
                            break

class Expression(Tag):
    def __init__(self,pattern):
            super().__init__(pattern)

    def if_condition(self, commands_dic, presentation, slide, shape, slides_index, dataObj):
        print("if>>>>>>>>>>>>>>>>>>>>>>>>>>>>>")
        pattern = CommandRegex.PATTERN_IF.value
        print("Pattern", pattern)
        matches = super().get_tag_content(pattern, shape)
        
        if( not matches or len(matches) < 1):
            return

        for match in matches:
            pattern_condition = r'\(\((.*?)\)\)'
            print("Match in matches", type(match))
            pattern_condition1 = CommandRegex.PATTERN_CONDITION.value[0]
            print("Pattern condition", type(pattern_condition))
            matched_condition = super().get_tag_from_string(pattern_condition,match)
            print('matched condition', matched_condition)

            pattern_content = CommandRegex.PATTERN_CONTENT.value[0]
            print("Pattern content", pattern_content)
            matched_content = super().get_tag_from_string(pattern_content,match)
            print("Mached_content", matched_content)
            
            for contidion in matched_condition:
                object_value = super().eval_executor(contidion, dataObj)

                #replace text
                text_replace_pattern = r'\+\+\+INS (.*?) \+\+\+'
                text_matches = super().get_tag_from_string(text_replace_pattern, matched_content[0])
                if( text_matches and len(text_matches) > 0):
                    text_replace = ""
                    if(object_value):
                        updated_data = super().text_tag_update(matched_content[0],dataObj)
                        if(updated_data and updated_data["text"]):
                            text_replace = updated_data["text"]
                    # this is not working if you use tabspaces, but you can use spaces
                    super().replace_tags(str(f"{CommandRegexSub.IF.value} {match}{CommandRegexSub.IF_END.value}"), text_replace, shape)

                #remove tables
                table_remove_pattern = CommandRegex.TABLE_REMOVE.value
                print("Table remove", table_remove_pattern)
                table_remove_matches = super().get_tag_from_string(table_remove_pattern, matched_content[0])
                print("table remove matches", table_remove_matches)
                if( table_remove_matches and len(table_remove_matches) > 0):
                    print("IF")
                    if(object_value):
                        print("Object value")
                        table = Table()
                        print("Slide", slide)
                        print("match content", matched_content)
                        table.remove_tables(slide,matched_content[0])
                        print("333333")
                        super().replace_tags(str(f"{CommandRegexSub.IF.value} {match}{CommandRegexSub.IF_END.value}"), "", shape)
                print("11111")

                #remove table row
                print("remove table row")
                table_row_remove_pattern = CommandRegex.TABLE_ROW_REMOVE.value
                table_row_remove_matches = super().get_tag_from_string(table_row_remove_pattern, matched_content[0])
                print("table row remove matches", table_row_remove_matches)
                if( table_row_remove_matches and len(table_row_remove_matches) > 0):
                    if(object_value):
                        print("IF ROW")
                        table = Table()
                        table.remove_table_rows(slide,matched_content[0])
                        print("After remove table row")
                        super().replace_tags(str(f"{CommandRegexSub.IF.value} {match}{CommandRegexSub.IF_END.value}"), "", shape)
                print("222222")        
                #remove table column
                table_column_remove_pattern = CommandRegex.TABLE_COLUMN_REMOVE.value[0]
                table_column_remove_matches = super().get_tag_from_string(table_column_remove_pattern, matched_content[0])
                if( table_column_remove_matches and len(table_column_remove_matches) > 0):
                    if(object_value):
                        Table.remove_table_column(slide,matched_content[0])
                        super().replace_tags(str(f"{CommandRegexSub.IF.value} {match}{CommandRegexSub.IF_END.value}"), "", shape)

    def for_loop(self, presentation, slide, shape, slides_index, dataObj):
        pattern = CommandRegex.PATTERN_FOR.value[0]
        matches = super().get_tag_content(pattern, shape)
        
        if( not matches or len(matches) < 1):
            return

        for match in matches:
            pattern_condition = CommandRegex.PATTERN_CONDITION.value
            matched_condition = super().get_tag_from_string(pattern_condition,match)

            pattern_content = CommandRegex.PATTERN_CONTENT.value
            matched_content = super().get_tag_from_string(pattern_content,match)
            
            for contidion in matched_condition:
                object_value = pydash.get(dataObj, contidion)
                text_result = ""
                if(object_value):
                    for data in object_value:
                        updated_data = super().text_tag_update(matched_content[0],data)
                        if(updated_data and updated_data["text"]):
                            text_result += updated_data["text"] + "\n"

            # this is not working if you use tabspaces, but you can use spaces
            super().replace_tags(str(f"+++FOR {match} FOR-END+++"), text_result, shape)
            

class CommandRegistry:
        def __init__(self):
            self.commands = {
                "+++IF": "if_condition",
                "+++FOR": "for_loop",
                "+++IM": "replace_image",
                "+++TB_ADD": "replace_table",
                "+++TB_TX_UP": "update_table_text",
                "+++TB_DRW": "draw_tables",
                "+++INS": "text_replace",
                "+++IF_INS": "if_condition",
                "+++TABLE_REMOVE": "if_condition"
            }

        def get_commands_dictionary(self):
            return self.commands

        def get_commands_list(self):
            return self.commands.keys()

        def get_command(self,command_name):
            if command_name == Command.IF_CONDITION:
                print("111")
                expression = Expression(Tag)
                print("222")
                return expression.if_condition
            elif command_name == Command.FOR_LOOP:
                expression = Expression()
                return expression.for_loop
            if command_name == Command.REPLACE_IMAGE:
              # def getData(shape,slide):
              #    return getData
              try:
                 image = Image()
                 def forTest(commands_dic,presentation,slide,shape,slides,dataObj):
                   
                    return image.replace_images(slide,r'\+\+\+IM (.*?) \+\+\+',shape)
                  #  print('aaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaa',a,b,c,d,e,f)
                 return forTest
              except Exception as e:
                print(e)
            elif command_name == Command.TEXT_REPLACE:
              # def getData(shape,slide):
              #    return getData
              try:
                 text = Text()
                 #text.text_replace()
                 def forTest(commands_dic,presentation,slide,shape,slides,dataObj):
                    return text.text_replace(slide,r'\+\+\+INS (.*?) \+\+\+',shape)
                 return forTest
              except Exception as e:
                print(e)         
            else:
                raise Exception("Invalid command name")


class CommandExecutor:
        registry = CommandRegistry()

        def __init__(self, presentation, dataObj):
            self.slide = None
            self.presentation = presentation
            self.dataObj = dataObj

        def execute(self):
            # get command names as a list
            commands = self.registry.get_commands_list()
            # print("a",commands)
            commands_dic = self.registry.get_commands_dictionary()
            # print("b",commands_dic)

            # find commands in the presentation using 'commands' list & execute
            slides = [self.slide for self.slide in self.presentation.slides]
            for shape in self.slide.shapes:
                if shape.has_text_frame:
                    print('slide', slides)
                    if shape.text:
                        try:
                            print("shape.text",shape.text)
                            # self.registry.get_command(Command.TEXT_REPLACE)(commands_dic, self.presentation, self.slide, shape, slides.index(self.slide), self.dataObj)
                            self.registry.get_command(Command.IF_CONDITION)(commands_dic, self.presentation, self.slide, shape, slides.index(self.slide), self.dataObj)
                        except Exception as e:
                           print(e.__class__)
                          #  print(commands_dic, self.presentation, self.slide, shape, slides.index(self.slide), self.dataObj)
                        
    

    
executor = CommandExecutor(prs, dataObj)
executor.execute()
prs.save('output.pptx')
    
  
    
   
        
        
