import os
from io import BytesIO
import boto3
from botocore.exceptions import ClientError
import logging
from pptx import Presentation
import json
import pydash
from tpip_pptx.constants import Command
from tpip_pptx.image import Image
from tpip_pptx.table import Table
from tpip_pptx.text import Text
from tpip_pptx.expression import Expression
from tpip_pptx.toc import Toc
from tpip_pptx.slide import Slide


POC_PPTX_BUCKET = os.environ.get("POC_PPTX_BUCKET")


class CommandRegistry:
    def __init__(self):
        self.commands = {
            "+++IF": "if_condition",
            "+++FOR": "for_loop",
            "+++IM": "replace_images",
            "+++TB_ADD": "replace_table",
            "+++TB_TX_UP": "update_table_text",
            "+++TB_DRW": "draw_tables",
            "+++INS": "text_replace"
        }

    def get_commands_dictionary(self):
        return self.commands

    def get_commands_list(self):
        return self.commands.keys()

    def get_command(self, command_name):
        if command_name == Command.REPLACE_IMAGE.value:
            image = Image()
            return image.replace_images
         
        elif command_name == Command.TEXT_REPLACE.value:
            text = Text()
            return text.text_replace
            
        elif command_name == Command.UPDATE_TABLE_TEXT.value:
            table = Table()
            return table.update_table_text
          
        elif command_name == Command.DRAW_TABLE.value:
            table = Table()
            return table.drow_tables

        elif command_name == Command.REPLACE_TABLE.value:
            table = Table()
            return table.replace_tables

        elif command_name == Command.REMOVE_TABLE.value:
            table = Table()
            return table.remove_tables

        elif command_name == Command.IF_CONDITION.value:
            expression = Expression()
            return expression.if_condition 

        elif command_name == Command.FOR_LOOP.value:
                expression = Expression()
                return expression.for_loop   
        else:
            raise Exception("Invalid command name")




class CommandExecutor:
    registry = CommandRegistry()

    def __init__(self, presentation, data_obj):
        self.presentation = presentation
        self.data_obj = data_obj

    def execute(self):
        # get command names as a list
        commands = self.registry.get_commands_list()

        commands_dic = self.registry.get_commands_dictionary()

        # find commands in the presentation using 'commands' list & execute
        slides = [self.slide for self.slide in self.presentation.slides]
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for cmd in commands:
                        if cmd in shape.text:
                            
                            try:
                                self.registry.get_command(commands_dic[cmd])(commands_dic, self.presentation, slide,
                                                                            shape, slides.index(slide), self.data_obj)
        
                            except Exception as e:
                                print(e)

        


def generate_pptx(event, context):
    json_data = json.loads(event['body'])
    data_obj = json.loads(pydash.get(json_data,'template'))
    template_key = pydash.get(json_data,'templateKey')
    dest_path = pydash.get(json_data,'destPath')
    output_file_name = pydash.get(json_data,'outPutFileName')
    dest_bucket_name = pydash.get(json_data,'destBucketName')
    bucket_path = str(f"{dest_path}/{output_file_name}")
    s3_client = boto3.client('s3')
    s3 = boto3.resource('s3')
    
    try:
        response = s3_client.get_object(Bucket=POC_PPTX_BUCKET, Key=str(f"{template_key}.pptx"))
        data = response['Body'].read()
        f = open(f'/tmp/{template_key}.pptx', "wb")
        f.write(data)
        f.close()
        presentation_object = Presentation(f'/tmp/{template_key}.pptx')
        executor = CommandExecutor(presentation_object, data_obj)
        executor.execute()
        slideObj = Slide()
        slideObj.remove_extra_slides(presentation_object)

        toc = Toc()
        toc.drow_toc(presentation_object, data_obj)
        

        try:
            with BytesIO() as fileobj:
                presentation_object.save(fileobj)
                fileobj.seek(0)
                bucket = s3.Bucket(output_file_name)
                s3.meta.client.head_bucket(Bucket=dest_bucket_name)
                res = s3_client.upload_fileobj(fileobj, dest_bucket_name, bucket_path)
                s3_bucket = boto3.resource("s3").Bucket(POC_PPTX_BUCKET)
                
                os.system("curl https://s3.amazonaws.com/lambda-libreoffice-demo/lo.tar.gz -o /tmp/lo.tar.gz && cd /tmp && file  /tmp/lo.tar.gz && cd /tmp && tar -xvf /tmp/lo.tar.gz")
                os.system('cd /tmp && file  /tmp/lo.tar.gz')
                convertCommand = "instdir/program/soffice --headless --invisible --nodefault --nofirststartwizard --nolockcheck --nologo --norestore --convert-to pdf --outdir /tmp"
                response = s3_client.get_object(Bucket=POC_PPTX_BUCKET, Key=bucket_path)
                data = response['Body'].read()
                f = open(f'/tmp/{output_file_name}', "wb")
                f.write(data)
                f.close()
                # Execute libreoffice to convert input file
                os.system(f"cd /tmp && {convertCommand} {output_file_name}")
                # Save converted object in S3
                outputFileName, _ = os.path.splitext(output_file_name)
                outputFileName = outputFileName  + ".pdf"
                f = open(f"/tmp/{outputFileName}","rb")
                s3_bucket.put_object(Key=str(f"{dest_path}/{outputFileName}"),Body=f)
                f.close()
            
                    
        except ClientError as e:
            logging.error(e.response)
            if e.response["Error"]["Code"] == "404":
                response={
                        "statusCode":404,
                        "message":"Destination bucket cannot be found"
                        }
                return response
            else:
                response={
                        "statusCode":500,
                        "message":"Something went wrong"
                        }
                return response
            

        body = {
            "message": "PPTX and PDF generate successfully!",
            "content": {
                "path":bucket_path
            }
        }

        response = {
            "statusCode": 200,
            "body": json.dumps(body)
        }

        return response
    
    
    except ClientError as e:
         logging.error(e.response)
         if e.response["Error"]["Code"] == "NoSuchKey":
            response= {
                "statusCode":404,
                "message":"The specified key does not exist."
            }
            return response
         else:
            response= {
                "statusCode":500,
                "message":"Something went wrong."
            }
            return response
